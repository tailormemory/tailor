"""
TAILOR — 4: Document Ingest (PDF, Excel, CSV, Word, PowerPoint, Markdown)
Scans configured folders (e.g. OneDrive) and ingests documents into the KB.

Features:
- PDF: extracts text per page via PyMuPDF
- Excel (.xlsx, .xls): extracts text per sheet via openpyxl
- CSV/TSV: extracts text as table
- Word (.docx): extracts text per paragraph + tables via python-docx
- PowerPoint (.pptx): extracts text per slide via python-pptx
- Markdown (.md, .markdown): extracts text per header section (header+fence aware)
- Chunking with v2 logic (target ~1200 chars)
- Incremental: tracks already-processed files (SHA256 hash)
- Source tag: "document" in metadata

Uso:
  python scripts/ingest/ingest_docs.py              # Incremental ingest
  python scripts/ingest/ingest_docs.py --full       # Re-ingest all documents
  python scripts/ingest/ingest_docs.py --list       # Show found files without ingesting
  python scripts/ingest/ingest_docs.py --status     # Show file status (new/modified/unchanged)

Requirements:
  pip install pymupdf openpyxl requests chromadb python-docx python-pptx

Configuration:
  Edit cloud_sync paths in config/tailor.yaml to point to desired folders.
"""

import csv
import datetime
import hashlib
import json
import os
import re
import sys
import time
import requests
import chromadb

# ============================================================
# CONFIGURAZIONE (from config/tailor.yaml)
# ============================================================
sys.path.insert(0, os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "lib"))
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__)))))
from embedding import get_embedding, get_embeddings
from embedding_contract import embedding_text
from config import get as cfg
from ingest_helpers import verified_upsert, make_run_id
from chunk_sidecar import invalidate_chunk_sidecars
from ocr_quality import assess_text_quality
from vlm_extractor import extract_page_via_vlm

# Cartelle da scansionare (ricorsivamente) — from YAML
WATCH_FOLDERS = cfg("ingest", "document_paths") or []

# Estensioni supportate
SUPPORTED_EXTENSIONS = set(cfg("ingest", "supported_extensions") or [".pdf", ".xlsx", ".xls", ".csv", ".tsv", ".docx", ".pptx", ".md", ".markdown"])

# Subfolders to ignore (exact names, case-insensitive)
IGNORE_FOLDERS = set(cfg("ingest", "ignore_folders") or ["Food Tracker"])

# Files to ignore (partial patterns in filename)
IGNORE_PATTERNS = set(cfg("ingest", "ignore_patterns") or ["~$", ".tmp", "Thumbs.db", ".DS_Store"])

DB_DIR = os.path.join(os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__)))), "db")
REGISTRY_FILE = os.path.join(DB_DIR, "doc_registry.json")
# Lock single-instance nel repo (non /tmp: multi-user, non encrypted).
LOCK_FILE = os.path.join(
    os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__)))),
    "ingest_docs.lock",
)

COLLECTION_NAME = cfg("kb", "collection") or "tailor_kb_v2"

BATCH_SIZE = 10

# Chunking config (aligned with 2_chunk.py v2)
TARGET_CHUNK_CHARS = 1200
MAX_CHUNK_CHARS = 2000
OVERLAP_CHARS = 200


# ============================================================
# FOLDER + DOC TYPE INFERENCE
# ============================================================

import re as _re


# # v1.2.4 paddle real impl
# Lazy module-level singleton for PPStructureV3.
# Init cost (warm cache): ~1.6s. Cold cache: triggers ~1GB model download from
# Paddle CDN (allocate 30-60min on first run, depending on bandwidth).
# Singleton avoids re-init across multiple PDFs in same process (e.g. nightly
# batch ingestion).
_paddle_pipeline = None


def _get_paddle_ocr():
    """Return lazy-initialized PPStructureV3 instance.

    Silences paddleocr/paddlex/ppocr loggers to WARNING to suppress per-predict
    'Creating model:' noise (cosmetic — models already in memory after init).
    """
    global _paddle_pipeline
    if _paddle_pipeline is None:
        import logging
        for _name in ("paddleocr", "paddlex", "ppocr"):
            logging.getLogger(_name).setLevel(logging.WARNING)
        from paddleocr import PPStructureV3
        _paddle_pipeline = PPStructureV3(
            lang="it",
            use_textline_orientation=True,
            use_table_recognition=True,
            use_seal_recognition=False,
            use_formula_recognition=False,
            use_chart_recognition=False,
            use_doc_orientation_classify=False,
            use_doc_unwarping=False,
        )
    return _paddle_pipeline


# Patterns that identify "bilancio gestionale centro-di-costo" templates
# (Sistemi/Zucchetti/TeamSystem). These have layout detection failures in
# PaddleOCR PPStructureV3 — see docs/v1.2.4_paddleocr_comparison.md.
# Detection: scan first page text via PyMuPDF; if any pattern matches, route
# to legacy extractor regardless of paddle config.
_GESTIONALE_PATTERNS = (
    "BILANCIO DI VERIFICA",
    "TCA PER CENTRO DI COSTO",
    "Centro di costo",
    "TCA per centro di costo",
)


def _is_bilancio_gestionale(filepath):
    """Return True if first page of PDF matches gestionale template patterns.

    Uses PyMuPDF quick text extract (no OCR). ~50-100ms cost per doc.
    Conservative: if extraction fails, returns False (paddle proceeds).
    """
    try:
        import fitz
        with fitz.open(filepath) as doc:
            if doc.page_count == 0:
                return False
            first_page_text = doc[0].get_text()
            for pat in _GESTIONALE_PATTERNS:
                if pat in first_page_text:
                    return True
    except Exception:
        return False
    return False


def _normalize_legacy_sections(sections):
    """Add 'markdown' key and 'metadata.extractor' to legacy sections so that
    whole-doc fallback returns a schema consistent with paddle output.
    Idempotent.
    """
    out = []
    for s in sections or []:
        if not isinstance(s, dict):
            out.append(s)
            continue
        meta = dict(s.get("metadata") or {})
        meta.setdefault("extractor", "legacy_fallback")
        out.append({
            "text": s.get("text", ""),
            "markdown": s.get("markdown", ""),
            "metadata": meta,
        })
    return out


def _log_paddle_fallback(filepath, reason, error, events):
    """Append a fallback audit record to logs/paddle_fallback_<YYYYMMDD>.json."""
    import json
    import logging
    from datetime import datetime
    from pathlib import Path as _Path

    logger = logging.getLogger(__name__)
    try:
        log_dir = _Path(os.environ.get("TAILOR_HOME", os.path.expanduser("~/tailor"))) / "logs"
        log_dir.mkdir(parents=True, exist_ok=True)
        log_file = log_dir / f"paddle_fallback_{datetime.now().strftime('%Y%m%d')}.json"

        entry = {
            "ts": datetime.now().isoformat(),
            "filepath": str(filepath),
            "reason": reason,
            "error": error,
            "events": events,
        }

        existing = []
        if log_file.exists():
            try:
                existing = json.loads(log_file.read_text())
                if not isinstance(existing, list):
                    existing = [existing]
            except Exception:
                existing = []

        existing.append(entry)
        log_file.write_text(json.dumps(existing, indent=2, ensure_ascii=False))
    except Exception as e:
        logger.error("Failed to write fallback audit log: %s", e)


# Folder → readable label map (built from cloud_sync config)
def _build_folder_labels():
    try:
        from config import get as _cfg
        labels = {}
        for sync in (_cfg("cloud_sync") or []):
            for folder in (sync.get("folders") or []):
                labels[folder] = folder
        return labels if labels else {}
    except Exception:
        return {}

FOLDER_LABELS = _build_folder_labels()

DOC_TYPE_PATTERNS = [
    # Financial statements and accounting
    (r"nota.integr",                    "nota_integrativa"),
    (r"bilancio.?(d['\s]esercizio|consolidat|abbreviat|annuale|ordinario)", "bilancio"),
    (r"stato.patrimoniale",             "bilancio"),
    (r"conto.economico",                "bilancio"),
    (r"xbrl",                           "bilancio"),
    (r"relazione.?(degli.amministratori|sulla.gestione)", "relazione_gestione"),
    (r"verbale.?(assemblea|cda|consiglio)", "verbale"),
    (r"statuto",                        "statuto"),
    # Fatture e pagamenti
    (r"fattur[ae]",                     "fattura"),
    (r"invoice",                        "fattura"),
    (r"ricevut[ae]",                    "ricevuta"),
    (r"bonifico",                       "pagamento"),
    # Estratti conto e investimenti
    (r"estratto.conto",                 "estratto_conto"),
    (r"rendiconto",                     "estratto_conto"),
    (r"portafoglio|portfolio",          "rendiconto_investimenti"),
    # Contratti e atti legali
    (r"atto.notarile|rogito",           "atto_notarile"),
    (r"contratto.?(di.?)?(locazione|affitto|compravendita|fornitura)", "contratto"),
    (r"scrittura.privata",              "contratto"),
    (r"visura",                         "visura_camerale"),
    (r"sentenza|decreto|ordinanza",     "atto_legale"),
    (r"ricorso|citazione",              "atto_legale"),
    # Tasse e fisco
    (r"dichiarazione.?(dei.redditi|iva|730|unico)", "dichiarazione_fiscale"),
    (r"\bf24\b|\bf24[_\-]|\bf23\b",     "pagamento_fiscale"),
    (r"unico\b",                        "dichiarazione_fiscale"),
    # Salute
    (r"refert[oi]|diagnos|anamnesi",    "referto_medico"),
    (r"esami?.?(del.sangue|clinici|laboratorio)", "referto_medico"),
    (r"cartella.clinica",               "referto_medico"),
    (r"ricetta",                        "prescrizione"),
    # Assicurazioni
    (r"polizza|assicuraz",              "polizza_assicurativa"),
    (r"sinistro|denuncia.sinistro",     "denuncia_sinistro"),
    # Immobili
    (r"catasto|planimetria",            "documento_catastale"),
    (r"mutuo|ipoteca",                  "mutuo"),
    # Presentazioni e piani
    (r"presentazione|pitch|deck|slide", "presentazione"),
    (r"piano.?(integrazione|strategico|operativo|marketing)", "piano"),
]

def infer_folder(filepath):
    """Estrae la cartella radice dal path del documento."""
    parts = filepath.replace("\\", "/").split("/")
    try:
        idx = parts.index("Personale")
        if idx + 1 < len(parts):
            candidate = parts[idx + 1]
            return FOLDER_LABELS.get(candidate, candidate)
    except ValueError:
        pass
    # Fallback: controlla direttamente i WATCH_FOLDERS
    for folder in WATCH_FOLDERS:
        if filepath.startswith(folder):
            return os.path.basename(folder)
    return ""

def infer_doc_type(filename, text_preview=""):
    """Inferisce il tipo documento dal nome file, con fallback sul testo."""
    combined = filename.lower()
    for pattern, doc_type in DOC_TYPE_PATTERNS:
        if _re.search(pattern, combined, _re.IGNORECASE):
            return doc_type
    # Fallback sul testo (primi 500 chars)
    if text_preview:
        preview = text_preview[:500].lower()
        for pattern, doc_type in DOC_TYPE_PATTERNS:
            if _re.search(pattern, preview, _re.IGNORECASE):
                return doc_type
    # Fallback per estensione
    ext = os.path.splitext(filename)[1].lower()
    if ext in (".xlsx", ".xls"):
        return "spreadsheet"
    if ext in (".csv", ".tsv"):
        return "csv"
    if ext == ".docx":
        return "documento_word"
    if ext == ".pptx":
        return "presentazione"
    return "documento"


# ============================================================
# FILE REGISTRY (tracking incrementale)
# ============================================================

def load_registry():
    """Load the registry of already-processed files."""
    if os.path.exists(REGISTRY_FILE):
        with open(REGISTRY_FILE, "r", encoding="utf-8") as f:
            return json.load(f)
    return {}


def mark_file_failed(registry, f, reason):
    """Registra un tentativo fallito senza mai spacciarlo per riuscito.

    L'hash conservato e' quello dell'ultima ingest RIUSCITA (o "" se non ce
    n'e' mai stata): registrare qui l'hash corrente farebbe classificare il
    file "unchanged" al giro dopo, quindi mai piu' ritentato.
    """
    existing = registry.get(f["filepath"], {})
    fail_count = existing.get("fail_count", 0) + 1
    registry[f["filepath"]] = {
        "hash": existing.get("hash", ""),
        "filename": f["filename"],
        "status": "failed",
        "fail_count": fail_count,
        "last_error": reason,
        "date_last_attempt": time.strftime("%Y-%m-%d %H:%M:%S"),
    }
    return fail_count


def acquire_single_instance_lock(lock_path=None):
    """flock non bloccante. Ritorna il file object (da tenere aperto) o None.

    Una seconda istanza NON attende: esce. Il lock e' rilasciato dall'OS alla
    chiusura del FD / uscita del processo, quindi non resta appeso dopo un
    crash come farebbe un lockfile a presenza.
    """
    import fcntl

    fd = open(lock_path or LOCK_FILE, "w")
    try:
        fcntl.flock(fd, fcntl.LOCK_EX | fcntl.LOCK_NB)
    except BlockingIOError:
        fd.close()
        return None
    fd.write(f"{os.getpid()}\n")
    fd.flush()
    return fd


def save_registry(registry):
    """Salva il registro."""
    os.makedirs(os.path.dirname(REGISTRY_FILE), exist_ok=True)
    with open(REGISTRY_FILE, "w", encoding="utf-8") as f:
        json.dump(registry, f, indent=2, ensure_ascii=False)


def file_hash(filepath):
    """SHA256 del file per detect modifiche."""
    h = hashlib.sha256()
    with open(filepath, "rb") as f:
        for block in iter(lambda: f.read(65536), b""):
            h.update(block)
    return h.hexdigest()


# ============================================================
# TEXT EXTRACTION
# ============================================================

def _ocr_vision(page):
    """OCR via Apple Vision framework (macOS)."""
    try:
        import Vision
        from Quartz import CGImageSourceCreateWithData, CGImageSourceCreateImageAtIndex
        from Foundation import NSData
        pix = page.get_pixmap(dpi=300)
        png = pix.tobytes("png")
        ns_data = NSData.dataWithBytes_length_(png, len(png))
        src = CGImageSourceCreateWithData(ns_data, None)
        if not src: return ""
        cg = CGImageSourceCreateImageAtIndex(src, 0, None)
        if not cg: return ""
        req = Vision.VNRecognizeTextRequest.alloc().init()
        req.setRecognitionLanguages_(["it", "en"])
        req.setRecognitionLevel_(1)
        req.setUsesLanguageCorrection_(True)
        handler = Vision.VNImageRequestHandler.alloc().initWithCGImage_options_(cg, None)
        success, err = handler.performRequests_error_([req], None)
        if not success: return ""
        results = req.results()
        if not results: return ""
        return "\n".join(r.topCandidates_(1)[0].string() for r in results)
    except Exception:
        return ""


def extract_pdf(filepath):
    """Estrae testo da PDF. Routing v1.2.4: PaddleOCR per doc strutturati o low-text, altrimenti pipeline nativa."""
    if cfg("paddleocr", "enabled", False):
        doc_type = infer_doc_type(os.path.basename(filepath))
        route, reason = _route_to_paddle(filepath, doc_type)
        if route:
            return _extract_pdf_paddle(filepath, doc_type, reason)
    return _extract_pdf_legacy(filepath)


FINANCIAL_DOCTYPES_VLM_FALLBACK = {
    "rendiconto_investimenti",
    "estratto_conto",
}


def _write_vlm_audit_log(filepath, doctype, total_pages, decisions, daily_count_after, daily_cap):
    from datetime import datetime
    log_dir = os.path.expanduser("~/tailor/logs")
    os.makedirs(log_dir, exist_ok=True)
    ts = datetime.now()
    log_path = os.path.join(log_dir, f"vlm_fallback_{ts.strftime('%Y%m%d_%H%M%S')}.json")
    payload = {
        "timestamp": ts.isoformat(),
        "filepath": filepath,
        "doctype": doctype,
        "total_pages": total_pages,
        "decisions": decisions,
        "vlm_calls_today_after": daily_count_after,
        "vlm_daily_cap": daily_cap,
    }
    try:
        with open(log_path, "w") as f:
            json.dump(payload, f, indent=2, ensure_ascii=False)
    except Exception as e:
        print(f"  [vlm] audit log write failed: {e}")


def extract_pdf_with_quality_fallback(filepath, doctype=None):
    """Run extract_pdf, then re-extract garbage sections via VLM for financial doctypes.

    Returns the (possibly modified) sections list. Always falls back to original
    sections on any VLM problem — never raises. Writes per-invocation audit log
    when at least one quality check is performed.
    """
    sections = extract_pdf(filepath)

    vlm_cfg = cfg("enrichment", "vlm_extraction") or {}
    enabled = vlm_cfg.get("enabled", False)
    daily_cap = vlm_cfg.get("daily_cap", 100)
    raster_dpi = vlm_cfg.get("raster_dpi", 144)

    if not enabled:
        return sections
    if not doctype or doctype not in FINANCIAL_DOCTYPES_VLM_FALLBACK:
        return sections
    if not sections:
        return sections

    from vlm_extractor import vlm_daily_count_get

    total_pages = sections[0].get("metadata", {}).get("total_pages", len(sections))
    decisions = []

    for i, section in enumerate(sections):
        text = section.get("text", "") or ""
        meta = section.get("metadata", {}) or {}
        page = meta.get("page", i + 1)
        quality = assess_text_quality(text)

        if not quality["is_garbage"]:
            if quality["char_count"] < 200:
                decisions.append({"page": page, "quality": quality,
                                  "action": "kept_original_edge_case_short"})
            else:
                decisions.append({"page": page, "quality": quality,
                                  "action": "kept_original_clean"})
            continue

        if vlm_daily_count_get() >= daily_cap:
            decisions.append({"page": page, "quality": quality,
                              "action": "vlm_skipped_daily_cap"})
            continue

        vlm_section = extract_page_via_vlm(filepath, page, total_pages,
                                           daily_cap=daily_cap,
                                           raster_dpi=raster_dpi)
        if vlm_section is None:
            decisions.append({"page": page, "quality": quality,
                              "action": "vlm_failed_all_backends"})
            continue

        sections[i] = vlm_section
        decisions.append({
            "page": page, "quality": quality, "action": "vlm_success",
            "extractor_after": vlm_section["metadata"]["extractor"],
        })

    _write_vlm_audit_log(filepath, doctype, total_pages, decisions,
                         vlm_daily_count_get(), daily_cap)
    return sections


def _route_to_paddle(filepath, doc_type):
    """True if PaddleOCR should handle this PDF, with reason ('doc_type'|'low_text'|None)."""
    # v1.2.4 gestionale guard
    if _is_bilancio_gestionale(filepath):
        return False, "gestionale_blocklist"

    structured = cfg("paddleocr", "structured_doc_types", []) or []
    if doc_type in structured:
        return True, "doc_type"

    min_text = cfg("paddleocr", "min_text_threshold", 100)
    try:
        import fitz
        with fitz.open(filepath) as probe:
            if sum(len(p.get_text()) for p in probe) < min_text:
                return True, "low_text"
    except Exception:
        pass  # Probe failure: legacy handler will retry with its own error handling
    return False, None


def _extract_pdf_paddle(filepath, doc_type, reason):
    """Extract PDF via PaddleOCR PPStructureV3 with per-page legacy fallback.

    Schema: returns list of sections, one per page (page number 1-indexed):
        [
            {
                "text": <flat text, backward-compat>,
                "markdown": <HTML/markdown with table structure, may be empty>,
                "metadata": {
                    "page": <1-indexed>,
                    "total_pages": <int>,
                    "ocr": True,
                    "extractor": "paddleocr" | "legacy_fallback" | ...,
                },
            },
            ...
        ]

    Per-page failure (paddle predict raises) -> fallback to legacy for that page,
    logged in logs/paddle_fallback_<ts>.json.
    Whole-doc failure (init or top-level predict raises) -> fallback to legacy
    for entire doc, normalized via _normalize_legacy_sections().
    """
    import logging
    logger = logging.getLogger(__name__)

    print(f"    [paddle] {os.path.basename(filepath)} (doc_type={doc_type or 'unknown'}; reason={reason})")

    fallback_events = []
    sections = []

    try:
        pipeline = _get_paddle_ocr()
    except Exception as e:
        logger.error("PaddleOCR init failed for %s: %s", filepath, e)
        _log_paddle_fallback(filepath, "init_failure", str(e), [])
        return _normalize_legacy_sections(_extract_pdf_legacy(filepath))

    try:
        pages = list(pipeline.predict(input=filepath))
    except Exception as e:
        logger.error("PaddleOCR predict failed for %s: %s", filepath, e)
        _log_paddle_fallback(filepath, "whole_doc_predict_failure", str(e), [])
        return _normalize_legacy_sections(_extract_pdf_legacy(filepath))

    total_pages = len(pages)
    if total_pages == 0:
        logger.warning("PaddleOCR returned 0 pages for %s, falling back", filepath)
        _log_paddle_fallback(filepath, "zero_pages", "no pages returned", [])
        return _normalize_legacy_sections(_extract_pdf_legacy(filepath))

    # For per-page fallback: extract legacy sections lazily, lookup by page number.
    legacy_sections_cache = None

    def _legacy_for_page(page_1indexed):
        nonlocal legacy_sections_cache
        if legacy_sections_cache is None:
            legacy_sections_cache = _extract_pdf_legacy(filepath)
        for ls in legacy_sections_cache:
            if isinstance(ls, dict) and (ls.get("metadata") or {}).get("page") == page_1indexed:
                return ls
        return None

    for idx, page_obj in enumerate(pages):
        page_1indexed = idx + 1
        try:
            md_text = ""
            if hasattr(page_obj, "markdown"):
                md = page_obj.markdown
                if isinstance(md, dict):
                    md_text = md.get("markdown_texts", "") or ""
                elif isinstance(md, str):
                    md_text = md

            flat_text = ""
            try:
                if "overall_ocr_res" in page_obj:
                    ocr_res = page_obj["overall_ocr_res"]
                    if ocr_res is not None and "rec_texts" in ocr_res:
                        rec_texts = ocr_res["rec_texts"]
                        if isinstance(rec_texts, (list, tuple)):
                            flat_text = "\n".join(str(t) for t in rec_texts)
            except Exception:
                pass

            if not flat_text and md_text:
                flat_text = re.sub(r"<[^>]+>", " ", md_text)
                flat_text = re.sub(r"\s+", " ", flat_text).strip()

            sections.append({
                "text": flat_text,
                "markdown": md_text,
                "metadata": {
                    "page": page_1indexed,
                    "total_pages": total_pages,
                    "ocr": True,
                    "extractor": "paddleocr",
                },
            })
        except Exception as e:
            logger.warning(
                "PaddleOCR page-level failure for %s page %d: %s, falling back",
                filepath, page_1indexed, e,
            )
            fallback_events.append({
                "page": page_1indexed,
                "error_type": type(e).__name__,
                "error": str(e),
            })
            try:
                leg = _legacy_for_page(page_1indexed)
                if leg is not None:
                    leg_meta = dict(leg.get("metadata") or {})
                    leg_meta["extractor"] = "legacy_fallback"
                    sections.append({
                        "text": leg.get("text", ""),
                        "markdown": "",
                        "metadata": leg_meta,
                    })
                else:
                    sections.append({
                        "text": "",
                        "markdown": "",
                        "metadata": {
                            "page": page_1indexed,
                            "total_pages": total_pages,
                            "ocr": True,
                            "extractor": "legacy_fallback_empty",
                        },
                    })
            except Exception as e2:
                logger.error("Legacy fallback also failed for %s page %d: %s",
                             filepath, page_1indexed, e2)
                sections.append({
                    "text": "",
                    "markdown": "",
                    "metadata": {
                        "page": page_1indexed,
                        "total_pages": total_pages,
                        "ocr": True,
                        "extractor": "double_failure",
                    },
                })

    if fallback_events:
        _log_paddle_fallback(filepath, "per_page_fallback", "see events", fallback_events)

    return sections

def _extract_pdf_legacy(filepath):
    """Pipeline pre-v1.2.4: testo nativo -> Tesseract -> Apple Vision."""
    import fitz
    sections = []
    ocr_used = False
    vision_used = False
    try:
        doc = fitz.open(filepath)
        for page_num in range(len(doc)):
            page = doc[page_num]
            text = page.get_text().strip()
            if not text or len(text) < 20:
                try:
                    tp = page.get_textpage_ocr(language="ita+eng", dpi=300, full=True)
                    tess_text = page.get_text(textpage=tp).strip()
                    if len(tess_text) > len(text):
                        text = tess_text
                        ocr_used = True
                except Exception: pass
            if not text or len(text) < 20:
                vision_text = _ocr_vision(page)
                if len(vision_text) > len(text):
                    text = vision_text
                    vision_used = True
            if text:
                sections.append({
                    "text": text,
                    "metadata": {"page": page_num + 1, "total_pages": len(doc), "ocr": ocr_used or vision_used}
                })
        doc.close()
        if vision_used:
            print(f"    (Apple Vision OCR applicato su alcune pagine)")
        elif ocr_used:
            print(f"    (Tesseract OCR applicato su alcune pagine)")
    except Exception as e:
        print(f"    ERRORE PDF {filepath}: {e}")
    return sections


# ---- Excel: header detection + serializzazione ------------------------------
#
# Il formato legacy (`a | b | c` per riga) perde l'associazione colonna→valore
# appena il chunk viene tagliato: una riga a metà chunk non ha più intestazione.
# Se il foglio ha un header riconoscibile serializziamo per riga
# ("<prima_cella>: Header2=val; Header3=val") ripetendo l'header su OGNI sezione.
#
# La detection è deliberatamente CONSERVATIVA: un falso positivo corrompe la
# semantica di tutto il foglio, un miss costa solo il formato legacy. Se non
# rileviamo un header il fallback è byte-identico al comportamento precedente.

_XL_MAX_HEADER_SCAN = 5      # righe iniziali candidate a essere l'header
_XL_MIN_HEADER_CELLS = 2     # celle testuali minime perché una riga sia header
_XL_SAMPLE_DATA_ROWS = 20    # righe dati campionate per validare l'header
_XL_OVERFLOW_TOLERANCE = 0   # colonne dati senza intestazione tollerate (0 = strict)

# "12,50", "1.234,56", "-3%", "€ 1.200", "2024-01-31", "31/12/24"
_XL_NUMERIC_RE = re.compile(r"^[+\-]?[\d\s.,']+\s*%?$")
_XL_DATEISH_RE = re.compile(r"^\d{1,4}\s*[-/.]\s*\d{1,2}\s*[-/.]\s*\d{1,4}")
_XL_CURRENCY = "€$£¥"


def _xl_is_empty(value):
    return value is None or (isinstance(value, str) and not value.strip())


def _xl_is_texty(value):
    """True solo per testo vero. Numeri, date, percentuali, valute e stringhe
    numeriche NON contano come testo (sono dati, non intestazioni)."""
    if _xl_is_empty(value):
        return False
    if isinstance(value, bool) or isinstance(value, (int, float)):
        return False
    if isinstance(value, (datetime.datetime, datetime.date, datetime.time)):
        return False
    if not isinstance(value, str):
        return False

    s = value.strip()
    stripped = s.strip(_XL_CURRENCY).strip()
    if _XL_NUMERIC_RE.match(stripped):
        return False
    if _XL_DATEISH_RE.match(s):
        return False
    # Serve almeno una lettera: "#1", "---", "(*)" non sono intestazioni.
    return any(ch.isalpha() for ch in s)


def _xl_fmt(value):
    """Rendering di una cella nel formato header-aware (non usato nel fallback)."""
    if value is None:
        return ""
    if isinstance(value, bool):
        return "true" if value else "false"
    if isinstance(value, datetime.datetime):
        if (value.hour, value.minute, value.second) == (0, 0, 0):
            return value.date().isoformat()
        return value.strftime("%Y-%m-%d %H:%M:%S")
    if isinstance(value, datetime.date):
        return value.isoformat()
    if isinstance(value, datetime.time):
        return value.strftime("%H:%M:%S")
    if isinstance(value, float):
        # xlrd restituisce ogni numero come float: 3.0 -> "3", non "3.0".
        if value.is_integer() and abs(value) < 1e15:
            return str(int(value))
        return repr(value)
    return str(value).strip()


def _xl_detect_header(rows):
    """Indice della riga header, o None. `rows`: griglia di valori tipizzati."""
    limit = min(_XL_MAX_HEADER_SCAN, len(rows))
    for i in range(limit):
        row = rows[i]
        non_empty = [(j, v) for j, v in enumerate(row) if not _xl_is_empty(v)]

        if len(non_empty) < _XL_MIN_HEADER_CELLS:
            # Riga vuota o riga-titolo (una sola cella): può precedere l'header.
            if len(non_empty) <= 1:
                continue
            return None

        texty = [(j, v) for j, v in non_empty if _xl_is_texty(v)]
        # Densa ma non testuale (es. dati numerici già da riga 0) -> nessun header.
        if len(texty) < _XL_MIN_HEADER_CELLS or len(texty) * 2 <= len(non_empty):
            return None

        hmax = max(j for j, _ in non_empty)
        data = [r for r in rows[i + 1: i + 1 + _XL_SAMPLE_DATA_ROWS]
                if any(not _xl_is_empty(v) for v in r)]
        if not data:
            return None

        # Larghezza: i dati non devono debordare oltre l'ultima colonna intestata
        # (oltre la tolleranza), altrimenti quella riga "header" copre solo una
        # parte del foglio.
        span = hmax + _XL_OVERFLOW_TOLERANCE
        for r in data:
            if any(not _xl_is_empty(v) for v in r[span + 1:]):
                return None

        # Foglio multi-blocco: se l'header si ripete più in basso (es. un P&L con
        # sezioni GROSS INCOME / EXPENSES / ASSETS, ognuna con la sua
        # intestazione) applicare il PRIMO header a tutto il foglio produce
        # spazzatura ("EXPENSES: Gennaio=Gennaio; ..."). Meglio il legacy.
        header_labels = {j: _xl_fmt(v).casefold() for j, v in non_empty}
        for r in rows[i + 1:]:
            repeats = sum(
                1 for j, label in header_labels.items()
                if j < len(r) and not _xl_is_empty(r[j]) and _xl_fmt(r[j]).casefold() == label
            )
            if repeats >= 2 and repeats * 2 >= len(header_labels):
                return None

        # Discriminante forte: almeno una colonna con header testuale e valori in
        # maggioranza NON testuali (numeri/date). Un foglio di prosa non ce l'ha.
        for j in range(hmax + 1):
            vals = [r[j] for r in data if j < len(r) and not _xl_is_empty(r[j])]
            if len(vals) < 2:
                continue
            if sum(1 for v in vals if not _xl_is_texty(v)) * 2 > len(vals):
                return i
        return None
    return None


def _xl_legacy_sections(sheet_name, raw_rows, total_sheets):
    """Formato storico — byte-identico al comportamento pre-header-detection."""
    rows = []
    for row in raw_rows:
        cells = [str(c) if c is not None else "" for c in row]
        if any(c.strip() for c in cells):
            rows.append(" | ".join(cells))
    if not rows:
        return []
    return [{
        "text": f"[Sheet: {sheet_name}]\n" + "\n".join(rows),
        "metadata": {"sheet": sheet_name, "total_sheets": total_sheets, "row_count": len(rows)},
    }]


def _xl_header_sections(sheet_name, rows, total_sheets, header_idx):
    """Sezioni pre-chunked: righe intere, header ripetuto su ogni blocco."""
    header_row = rows[header_idx]
    hmax = max(j for j, v in enumerate(header_row) if not _xl_is_empty(v))
    # Colonne dati senza intestazione (tipicamente un totale in coda): entro la
    # tolleranza vengono emesse con etichetta posizionale invece che perse.
    base = hmax
    for r in rows[header_idx + 1:]:
        for j in range(base + 1, min(len(r), base + 1 + _XL_OVERFLOW_TOLERANCE)):
            if not _xl_is_empty(r[j]):
                hmax = max(hmax, j)
    labels = [
        _xl_fmt(header_row[j]) if j < len(header_row) and not _xl_is_empty(header_row[j])
        else f"Col{j + 1}"
        for j in range(hmax + 1)
    ]
    prefix = f"[Sheet: {sheet_name}]\nHeader: " + " | ".join(labels)

    lines = []
    for row in rows[header_idx + 1:]:
        if all(_xl_is_empty(v) for v in row):
            continue
        first = _xl_fmt(row[0]) if len(row) > 0 else ""
        parts = [
            f"{labels[j]}={_xl_fmt(row[j])}"
            for j in range(1, min(hmax + 1, len(row)))
            if not _xl_is_empty(row[j])
        ]
        if first and parts:
            lines.append(f"{first}: " + "; ".join(parts))
        elif first:
            lines.append(first)
        elif parts:
            lines.append("; ".join(parts))
    if not lines:
        return []

    # Blocchi di righe INTERE sotto TARGET_CHUNK_CHARS: il chunker generico li
    # passa as-is (nessuno split a metà riga, nessun overlap).
    sections = []
    block, block_len = [], len(prefix)
    for line in lines:
        if block and block_len + 1 + len(line) > TARGET_CHUNK_CHARS:
            sections.append((block, block_len))
            block, block_len = [], len(prefix)
        block.append(line)
        block_len += 1 + len(line)
    if block:
        sections.append((block, block_len))

    return [{
        "text": prefix + "\n" + "\n".join(b),
        "metadata": {
            "sheet": sheet_name,
            "total_sheets": total_sheets,
            "row_count": len(b),
            "prechunked": True,
        },
    } for b, _ in sections]


def _xl_sheet_sections(sheet_name, raw_rows, typed_rows, total_sheets):
    """Header-aware se rilevabile, altrimenti fallback legacy byte-identico."""
    header_idx = _xl_detect_header(typed_rows)
    if header_idx is not None:
        sections = _xl_header_sections(sheet_name, typed_rows, total_sheets, header_idx)
        if sections:
            return sections
    return _xl_legacy_sections(sheet_name, raw_rows, total_sheets)


def extract_excel(filepath):
    """Estrae testo da Excel, una sezione per sheet. Supporta .xlsx e .xls."""
    ext = os.path.splitext(filepath)[1].lower()
    if ext == ".xls":
        return _extract_xls(filepath)
    else:
        return _extract_xlsx(filepath)


def _extract_xlsx(filepath):
    """Estrae testo da .xlsx via openpyxl."""
    from openpyxl import load_workbook

    sections = []
    try:
        wb = load_workbook(filepath, read_only=True, data_only=True)
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            # openpyxl (data_only) restituisce già valori nativi: raw e typed
            # coincidono, il fallback resta byte-identico.
            rows = [list(row) for row in ws.iter_rows(values_only=True)]
            sections.extend(_xl_sheet_sections(sheet_name, rows, rows, len(wb.sheetnames)))
        wb.close()
    except Exception as e:
        print(f"    ERRORE Excel {filepath}: {e}")
    return sections


def _extract_xls(filepath):
    """Estrae testo da .xls via xlrd."""
    import xlrd

    sections = []
    try:
        wb = xlrd.open_workbook(filepath)
        for sheet_name in wb.sheet_names():
            ws = wb.sheet_by_name(sheet_name)
            raw_rows, typed_rows = [], []
            for row_idx in range(ws.nrows):
                raw, typed = [], []
                for col_idx in range(ws.ncols):
                    cell = ws.cell(row_idx, col_idx)
                    # raw = valore grezzo xlrd (il fallback fa str() come prima)
                    raw.append(cell.value)
                    typed.append(_xls_typed_value(cell, wb.datemode))
                raw_rows.append(raw)
                typed_rows.append(typed)
            sections.extend(_xl_sheet_sections(sheet_name, raw_rows, typed_rows, wb.nsheets))
    except Exception as e:
        print(f"    ERRORE Excel {filepath}: {e}")
    return sections


def _xls_typed_value(cell, datemode):
    """xlrd non tipizza: date e bool arrivano come float. Normalizza per la
    detection (il path legacy continua a usare il valore grezzo)."""
    import xlrd

    if cell.ctype in (xlrd.XL_CELL_EMPTY, xlrd.XL_CELL_BLANK, xlrd.XL_CELL_ERROR):
        return None
    if cell.ctype == xlrd.XL_CELL_BOOLEAN:
        return bool(cell.value)
    if cell.ctype == xlrd.XL_CELL_DATE:
        try:
            return xlrd.xldate_as_datetime(cell.value, datemode)
        except Exception:
            return float(cell.value)
    return cell.value


def extract_csv(filepath):
    """Estrae testo da CSV/TSV."""
    sections = []
    try:
        ext = os.path.splitext(filepath)[1].lower()
        delimiter = "\t" if ext == ".tsv" else ","

        with open(filepath, "r", encoding="utf-8", errors="replace") as f:
            reader = csv.reader(f, delimiter=delimiter)
            rows = []
            for row in reader:
                cells = [str(c).strip() for c in row]
                if any(c for c in cells):
                    rows.append(" | ".join(cells))

        if rows:
            text = "\n".join(rows)
            sections.append({
                "text": text,
                "metadata": {"row_count": len(rows)}
            })
    except Exception as e:
        print(f"    ERRORE CSV {filepath}: {e}")
    return sections


def extract_docx(filepath):
    """Estrae testo da .docx via python-docx. Paragrafi + tabelle in ordine."""
    from docx import Document

    sections = []
    try:
        doc = Document(filepath)
        parts = []

        # Itera sugli elementi del body in ordine di apparizione
        for element in doc.element.body:
            tag = element.tag.split("}")[-1] if "}" in element.tag else element.tag

            if tag == "p":
                text = element.text
                if text and text.strip():
                    parts.append(text.strip())

            elif tag == "tbl":
                for table in doc.tables:
                    if table._element is element:
                        table_rows = []
                        for row in table.rows:
                            cells = [cell.text.strip() for cell in row.cells]
                            if any(c for c in cells):
                                table_rows.append(" | ".join(cells))
                        if table_rows:
                            parts.append("[Tabella]\n" + "\n".join(table_rows))
                        break

        if parts:
            full_text = "\n\n".join(parts)
            sections.append({
                "text": full_text,
                "metadata": {
                    "paragraph_count": len([p for p in doc.paragraphs if p.text.strip()]),
                    "table_count": len(doc.tables),
                }
            })

    except Exception as e:
        print(f"    ERRORE DOCX {filepath}: {e}")
    return sections


def extract_pptx(filepath):
    """Estrae testo da .pptx via python-pptx. Una sezione per slide."""
    from pptx import Presentation

    sections = []
    try:
        prs = Presentation(filepath)
        total_slides = len(prs.slides)

        for slide_num, slide in enumerate(prs.slides, 1):
            parts = []

            # Titolo
            if slide.shapes.title and slide.shapes.title.has_text_frame:
                title_text = slide.shapes.title.text.strip()
                if title_text:
                    parts.append(f"[Titolo] {title_text}")

            # Contenuto shape (escluso titolo)
            for shape in slide.shapes:
                if shape == slide.shapes.title:
                    continue

                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        p_text = paragraph.text.strip()
                        if p_text:
                            parts.append(p_text)

                # Tabelle nelle slide
                if shape.has_table:
                    table_rows = []
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        if any(c for c in cells):
                            table_rows.append(" | ".join(cells))
                    if table_rows:
                        parts.append("[Tabella]\n" + "\n".join(table_rows))

            # Note speaker
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    parts.append(f"[Note] {notes_text}")

            if parts:
                text = f"[Slide {slide_num}/{total_slides}]\n" + "\n".join(parts)
                sections.append({
                    "text": text,
                    "metadata": {
                        "slide_number": slide_num,
                        "total_slides": total_slides,
                    }
                })

    except Exception as e:
        print(f"    ERRORE PPTX {filepath}: {e}")
    return sections


def _split_markdown_by_headers(text):
    """Segmenta il markdown in blocchi (heading, testo) a ogni header H2/H3,
    ignorando i '#' dentro i code fence (```/~~~). Il preambolo prima del primo
    header resta nel primo blocco con heading="". Header-aware: così il chunker
    paragraph-first non stacca un'intestazione dal suo corpo.
    """
    import re as _re
    lines = text.splitlines(keepends=True)
    blocks = []
    cur_heading = ""
    cur = []
    fence_char = None  # '`' o '~' se dentro un code fence, altrimenti None

    for line in lines:
        s = line.lstrip()
        if s[:3] in ("```", "~~~"):
            ch = s[0]
            fence_char = ch if fence_char is None else (None if ch == fence_char else fence_char)
            cur.append(line)
            continue
        if fence_char is None and _re.match(r"#{2,3} ", line):
            joined = "".join(cur).strip()
            if joined:
                blocks.append((cur_heading, joined))
            cur_heading = line.strip().lstrip("#").strip()
            cur = [line]
        else:
            cur.append(line)

    joined = "".join(cur).strip()
    if joined:
        blocks.append((cur_heading, joined))
    return blocks


def extract_markdown(filepath):
    """Estrae testo da Markdown (.md/.markdown), una sezione per header H2/H3.

    Header-aware + fence-aware: ogni sezione parte dalla propria intestazione e i
    '#' dentro i code fence non vengono scambiati per header. Le sezioni oltre il
    target vengono spezzate da split_markdown_at_boundaries() (path dedicato ai
    .md): code fence atomici, tabelle spezzate solo a confine di riga, prosa a
    confine di frase.
    """
    sections = []
    try:
        with open(filepath, "r", encoding="utf-8", errors="replace") as f:
            text = f.read()
    except Exception as e:
        print(f"    ERRORE Markdown {filepath}: {e}")
        return sections

    if not text.strip():
        return sections

    blocks = _split_markdown_by_headers(text)
    total = len(blocks)
    for i, (heading, body) in enumerate(blocks, 1):
        sections.append({
            "text": body,
            "metadata": {
                "section": i,
                "total_sections": total,
                "heading": heading,
            },
        })
    return sections


def extract_text(filepath, doctype=None):
    """Router: estrae testo in base all'estensione.

    doctype, if provided, routes PDF extraction through the VLM quality-fallback
    wrapper. Other formats ignore it (no OCR garbage risk on native text formats).
    """
    ext = os.path.splitext(filepath)[1].lower()
    if ext == ".pdf":
        return extract_pdf_with_quality_fallback(filepath, doctype=doctype)
    elif ext in (".xlsx", ".xls"):
        return extract_excel(filepath)
    elif ext in (".csv", ".tsv"):
        return extract_csv(filepath)
    elif ext == ".docx":
        return extract_docx(filepath)
    elif ext == ".pptx":
        return extract_pptx(filepath)
    elif ext in (".md", ".markdown"):
        return extract_markdown(filepath)
    return []


# ============================================================
# CHUNKING v3 (semantic paragraph-first)
# ============================================================

def split_text_at_boundaries(text, max_chars):
    """Spezza un testo lungo in blocchi a confini naturali — v3 (paragraph-first).
    
    Strategia: splitta per paragrafi, raggruppa fino al target, merge frammenti piccoli.
    Rispetta confini semantici naturali invece di tagliare a max_chars e cercare all'indietro.
    """
    import re as _re
    MIN_CHUNK = 300
    
    if len(text) <= max_chars:
        return [text]
    
    # Step 1: split per paragrafi (doppio newline)
    paragraphs = _re.split(r'\n\s*\n', text)
    paragraphs = [p.strip() for p in paragraphs if p.strip()]
    
    if not paragraphs:
        return [text[:max_chars]]
    
    # Step 2: raggruppa paragrafi fino al target
    chunks = []
    current_parts = []
    current_chars = 0
    
    for para in paragraphs:
        para_len = len(para)
        
        # Paragrafo singolo troppo lungo -> spezza a frase
        if para_len > max_chars:
            if current_parts:
                chunks.append('\n\n'.join(current_parts))
                current_parts = []
                current_chars = 0
            remaining = para
            while len(remaining) > max_chars:
                cut = -1
                seg = remaining[:max_chars]
                for sep in ['. ', '.\n', '? ', '! ', '; ']:
                    pos = seg.rfind(sep)
                    if pos > max_chars * 0.3:
                        cut = pos + len(sep)
                        break
                if cut == -1:
                    pos = seg.rfind(' ')
                    cut = pos if pos > max_chars * 0.3 else max_chars
                chunks.append(remaining[:cut].strip())
                remaining = remaining[cut:].strip()
            if remaining:
                current_parts = [remaining]
                current_chars = len(remaining)
            continue
        
        would_be = current_chars + para_len + (2 if current_parts else 0)
        
        if would_be > max_chars and current_chars >= MIN_CHUNK:
            chunks.append('\n\n'.join(current_parts))
            current_parts = [para]
            current_chars = para_len
        else:
            current_parts.append(para)
            current_chars = would_be
    
    if current_parts:
        chunks.append('\n\n'.join(current_parts))
    
    # Step 3: post-hoc merge — chunks too small get merged
    merged = []
    i = 0
    while i < len(chunks):
        chunk = chunks[i]
        while len(chunk) < MIN_CHUNK and i + 1 < len(chunks):
            i += 1
            chunk = chunk + '\n\n' + chunks[i]
        if len(chunk) < MIN_CHUNK and merged:
            merged[-1] = merged[-1] + '\n\n' + chunk
        else:
            merged.append(chunk)
        i += 1
    
    return merged


# ---- Markdown-aware chunking (path .md only — NON tocca split_text_at_boundaries) ----

def _is_md_table_block(block):
    """True se il blocco è una tabella markdown (maggioranza righe con '|')."""
    lines = [l for l in block.splitlines() if l.strip()]
    if len(lines) < 2:
        return False
    return sum(1 for l in lines if "|" in l) >= len(lines) * 0.6


def _segment_markdown(text):
    """Spezza una sezione markdown in segmenti (kind, testo).

    kind='fence': code block ```/~~~ completo -> ATOMICO, mai spezzato.
    kind='para' : paragrafo (blank-line delimited) fuori dai fence.
    Un fence non terminato a fine sezione resta comunque atomico (il source è
    già sbilanciato: non c'è nulla da bilanciare).
    """
    import re as _re
    lines = text.splitlines(keepends=True)
    segments = []
    buf = []
    fence_char = None
    fence_buf = []

    def _flush_buf():
        if buf:
            for p in _re.split(r'\n\s*\n', "".join(buf)):
                p = p.strip()
                if p:
                    segments.append(("para", p))
            buf.clear()

    for line in lines:
        s = line.lstrip()
        if fence_char is None and s[:3] in ("```", "~~~"):
            _flush_buf()
            fence_char = s[0]
            fence_buf = [line]
            continue
        if fence_char is not None:
            fence_buf.append(line)
            if s.startswith(fence_char * 3):
                segments.append(("fence", "".join(fence_buf).strip("\n")))
                fence_char = None
                fence_buf = []
            continue
        buf.append(line)

    if fence_char is not None:           # fence non chiuso -> atomico comunque
        segments.append(("fence", "".join(fence_buf).strip("\n")))
    _flush_buf()
    return segments


def _split_md_table_rows(table, max_chars):
    """Raggruppa le righe di una tabella fino a max_chars, mai a metà riga.
    Una riga singola più lunga del max resta intera (oversize ma non spezzata)."""
    rows = table.splitlines()
    out, cur, cur_len = [], [], 0
    for row in rows:
        rlen = len(row) + 1
        if cur and cur_len + rlen > max_chars:
            out.append("\n".join(cur))
            cur, cur_len = [row], rlen
        else:
            cur.append(row)
            cur_len += rlen
    if cur:
        out.append("\n".join(cur))
    return out


def _split_md_plain(text, max_chars):
    """Spezza prosa lunga a confine di frase/spazio (nessuna struttura da preservare)."""
    out = []
    remaining = text
    while len(remaining) > max_chars:
        seg = remaining[:max_chars]
        cut = -1
        for sep in ('. ', '.\n', '? ', '! ', '; ', '\n'):
            pos = seg.rfind(sep)
            if pos > max_chars * 0.3:
                cut = pos + len(sep)
                break
        if cut == -1:
            pos = seg.rfind(' ')
            cut = pos if pos > max_chars * 0.3 else max_chars
        out.append(remaining[:cut].strip())
        remaining = remaining[cut:].strip()
    if remaining:
        out.append(remaining)
    return out


def split_markdown_at_boundaries(text, max_chars):
    """Splitter markdown-aware per le sole sezioni .md.

    Garanzie sul testo prodotto:
      - code fence (```/~~~) bilanciati: blocco atomico, se supera max_chars
        diventa un chunk a sé (oversize) ma mai spezzato internamente;
      - tabelle spezzate solo a confine di riga, mai a metà riga;
      - prosa lunga spezzata a confine di frase/spazio.
    Ogni chunk è bilanciato per costruzione (i fence sono segmenti interi; i
    chunk para/tabella/prosa non contengono marker di fence), quindi anche il
    merge per concatenazione resta bilanciato.
    Stesso contratto di split_text_at_boundaries (lista di stringhe). NON usata
    per PDF/docx/pptx -> comportamento di quei formati invariato.
    """
    MIN_CHUNK = 300
    if len(text) <= max_chars:
        return [text]

    segments = _segment_markdown(text)
    chunks = []
    cur = []
    cur_chars = 0

    def _flush():
        nonlocal cur_chars
        if cur:
            chunks.append("\n\n".join(cur))
            cur.clear()
            cur_chars = 0

    for kind, seg in segments:
        if not seg:
            continue
        seg_len = len(seg)

        # Segmento che da solo sfora il max
        if seg_len > max_chars:
            _flush()
            if kind == "fence":
                chunks.append(seg)                        # atomico -> mai spezzato
            elif _is_md_table_block(seg):
                chunks.extend(_split_md_table_rows(seg, max_chars))
            else:
                chunks.extend(_split_md_plain(seg, max_chars))
            continue

        would_be = cur_chars + seg_len + (2 if cur else 0)
        if would_be > max_chars and cur_chars >= MIN_CHUNK:
            _flush()
            cur.append(seg)
            cur_chars = seg_len
        else:
            cur.append(seg)
            cur_chars = would_be

    _flush()

    # Post-hoc merge dei chunk piccoli (concat \n\n -> preserva bilanciamento
    # fence e confini di riga; mai split interno).
    merged = []
    i = 0
    while i < len(chunks):
        chunk = chunks[i]
        while len(chunk) < MIN_CHUNK and i + 1 < len(chunks):
            i += 1
            chunk = chunk + '\n\n' + chunks[i]
        if len(chunk) < MIN_CHUNK and merged:
            merged[-1] = merged[-1] + '\n\n' + chunk
        else:
            merged.append(chunk)
        i += 1
    return merged


def chunk_document(sections, file_info):
    """Converte le sezioni estratte in chunk per ChromaDB."""
    chunks = []
    chunk_index = 0
    prev_chunk_text = ""

    # Compute folder and doc_type once per file
    folder = infer_folder(file_info["filepath"])
    text_preview = sections[0]["text"] if sections else ""
    doc_type = infer_doc_type(file_info["filename"], text_preview)

    ext_lower = (file_info.get("extension", "") or "").lower()
    is_md = ext_lower in (".md", ".markdown")

    for section in sections:
        text = section["text"]
        sec_meta = dict(section["metadata"])

        # Sezioni già tagliate a monte su confini semantici (fogli Excel
        # header-aware): passano as-is, senza split generico e senza overlap —
        # entrambi spezzerebbero righe e duplicherebbero l'header.
        prechunked = bool(sec_meta.pop("prechunked", False))

        # Spezza sezioni lunghe — markdown-aware per i .md (fence/tabelle preservati),
        # splitter condiviso (invariato) per tutti gli altri formati.
        if prechunked:
            parts = [text]
        elif is_md:
            parts = split_markdown_at_boundaries(text, TARGET_CHUNK_CHARS)
        else:
            parts = split_text_at_boundaries(text, TARGET_CHUNK_CHARS)

        for part in parts:
            # Aggiungi overlap — disattivato per .md: lo slice raw di OVERLAP_CHARS
            # introdurrebbe fence sbilanciati o frammenti di riga-tabella nel prefisso.
            overlap_text = ""
            if prev_chunk_text and OVERLAP_CHARS > 0 and not is_md and not prechunked:
                overlap_text = prev_chunk_text[-OVERLAP_CHARS:].strip()
                space_pos = overlap_text.find(" ")
                if space_pos > 0:
                    overlap_text = overlap_text[space_pos + 1:]

            if overlap_text:
                full_text = f"[...contesto precedente:] {overlap_text}\n\n{part}"
            else:
                full_text = part

            # Genera ID unico
            file_hash_short = file_info["hash"][:12]
            chunk_id = f"doc_{file_hash_short}_chunk_{chunk_index:04d}"

            # Metadata
            meta = {
                "conv_id": f"doc_{file_hash_short}",
                "title": file_info["filename"],
                "date": file_info["modified_date"],
                "create_time": file_info["modified_ts"],
                "default_model": "",
                "chunk_index": chunk_index,
                "char_count": len(full_text),
                "turn_count": 1,
                "source": "document",
                "category": file_info.get("category", ""),
                "folder": folder,
                "doc_type": doc_type,
                "file_type": file_info["extension"],
                "file_path": file_info["rel_path"],
            }
            # Aggiungi metadata sezione-specifici
            meta.update({f"sec_{k}": v for k, v in sec_meta.items()})

            chunks.append({
                "chunk_id": chunk_id,
                "text": full_text,
                "metadata": meta
            })

            prev_chunk_text = part
            chunk_index += 1

    return chunks


# ============================================================
# EMBEDDING + INGEST
# ============================================================



def embed_chunks(chunks):
    """Calcola gli embedding PRIMA che la KB venga toccata.

    Separata da ingest_chunks per tenere il cleanup il piu' vicino possibile
    all'upsert: se l'embedding fallisce non abbiamo ancora cancellato nulla e
    il documento vecchio resta intatto. Ritorna (batches, errors).
    """
    batches = []
    errors = 0

    for batch_start in range(0, len(chunks), BATCH_SIZE):
        batch = chunks[batch_start:batch_start + BATCH_SIZE]

        # Testo embeddato == testo salvato: unica fonte via contratto.
        documents = [embedding_text("document", c["metadata"], c["text"]) for c in batch]

        # Batch fallito -> contato in errors, gli altri proseguono.
        try:
            embeddings = get_embeddings(documents)
        except Exception as e:
            print(f"\n  Embedding error: {e}")
            errors += len(batch)
            continue

        batches.append({
            "ids": [c["chunk_id"] for c in batch],
            "embeddings": embeddings,
            "documents": documents,
            "metadatas": [c["metadata"] for c in batch],
        })

    return batches, errors


def ingest_chunks(collection, batches, run_id: str | None = None):
    """Scrive in ChromaDB batch gia' embeddati da embed_chunks()."""
    processed = 0
    errors = 0
    run_id = run_id or make_run_id("document")

    for b in batches:
        size = len(b["ids"])
        try:
            ok = verified_upsert(
                collection, b["ids"], b["embeddings"], b["documents"], b["metadatas"],
                run_id=run_id, source="document",
            )
            if not ok:
                # Non contarlo anche in processed: e' il criterio con cui il
                # chiamante decide se il registry puo' registrare successo.
                errors += size
                continue
        except Exception as e:
            print(f"\n  ChromaDB errore: {e}")
            errors += size
            continue

        processed += size

    return processed, errors


def delete_file_chunks(collection, file_hash_short):
    """Cancella i chunk di un file specifico dalla collection.

    FAIL-LOUD di proposito: l'eccezione propaga. Chi chiama deve saltare il
    file, perche' fare upsert dopo una delete non verificata lascia in KB un
    mix di chunk vecchi e nuovi con gli stessi id.
    """
    results = collection.get(
        where={"conv_id": f"doc_{file_hash_short}"},
        include=[]
    )
    if results["ids"]:
        collection.delete(ids=results["ids"])
        return len(results["ids"])
    return 0


def _pre_upsert_cleanup(collection, f, registry):
    """Ripulisce KB e sidecar prima di (ri)scrivere i chunk di un file.

    Vale per OGNI file processato, non solo i modificati: i chunk_id sono
    deterministici sull'hash del contenuto, quindi con `--full` un file a
    contenuto invariato riusa gli stessi id. Se nel frattempo cambiano i
    boundary di chunking (estrattore nuovo) il testo sotto quegli id cambia
    senza che nulla se ne accorga: extraction_log farebbe da gate di skip e il
    chunk non verrebbe mai piu' ri-estratto dal nightly.

    Solleva se la delete Chroma o l'invalidazione sidecar fallisce: il file va
    saltato, mai upsertato a meta'.
    """
    current = f["hash"][:12]
    hashes = [current]

    old_hash = (registry.get(f["filepath"], {}).get("hash", "") or "")[:12]
    if old_hash and old_hash != current:
        hashes.append(old_hash)

    for h in hashes:
        deleted = delete_file_chunks(collection, h)
        if deleted:
            print(f"  Rimossi {deleted} chunk vecchi (doc_{h})")
        invalidate_chunk_sidecars(f"doc_{h}", verbose=True)


# ============================================================
# FILE DISCOVERY
# ============================================================

def scan_folders():
    """Scansiona le cartelle configurate e restituisce i file supportati."""
    files = []
    for folder in WATCH_FOLDERS:
        if not os.path.exists(folder):
            print(f"  WARNING: folder not found: {folder}")
            continue

        for root, dirs, filenames in os.walk(folder):
            # Escludi sottocartelle in IGNORE_FOLDERS (modifica dirs in-place)
            dirs[:] = [d for d in dirs if d not in IGNORE_FOLDERS]

            for fname in filenames:
                # Ignore temporary files
                if any(p in fname for p in IGNORE_PATTERNS):
                    continue

                ext = os.path.splitext(fname)[1].lower()
                if ext not in SUPPORTED_EXTENSIONS:
                    continue

                filepath = os.path.join(root, fname)
                rel_path = os.path.relpath(filepath, os.path.commonpath(WATCH_FOLDERS + [filepath]))

                # Category from first subfolder
                parts = os.path.relpath(filepath, folder).split(os.sep)
                category = parts[0] if len(parts) > 1 else ""

                try:
                    stat = os.stat(filepath)
                    mod_ts = stat.st_mtime
                    mod_date = time.strftime("%Y-%m-%d", time.localtime(mod_ts))
                    size_kb = stat.st_size / 1024
                except Exception:
                    mod_ts = 0
                    mod_date = ""
                    size_kb = 0

                files.append({
                    "filepath": filepath,
                    "filename": fname,
                    "rel_path": rel_path,
                    "extension": ext,
                    "category": category,
                    "modified_ts": mod_ts,
                    "modified_date": mod_date,
                    "size_kb": size_kb,
                })

    return files


def classify_files(files, registry, full_mode=False):
    """Classifica file in: new, modified, unchanged. Deduplica per hash."""
    new_files = []
    modified_files = []
    unchanged_files = []
    seen_hashes = set()
    duplicates = 0
    permanently_failed = 0

    for f in files:
        fhash = file_hash(f["filepath"])
        f["hash"] = fhash

        # Dedup: skip if same content already seen
        if fhash in seen_hashes:
            duplicates += 1
            continue
        seen_hashes.add(fhash)

        # Skip permanently failed files (3+ extraction failures)
        reg_entry = registry.get(f["filepath"], {})
        if reg_entry.get("status") == "failed" and reg_entry.get("fail_count", 0) >= 3:
            permanently_failed += 1
            print(f"  [SKIP] {f['filename']} — permanently failed after {reg_entry['fail_count']} attempts")
            continue

        if full_mode:
            new_files.append(f)
        elif f["filepath"] not in registry:
            new_files.append(f)
        elif registry[f["filepath"]]["hash"] != fhash:
            modified_files.append(f)
        else:
            unchanged_files.append(f)

    if duplicates > 0:
        print(f"  Duplicates skipped (same content): {duplicates}")
    if permanently_failed > 0:
        print(f"  Permanently failed (skipped): {permanently_failed}")

    return new_files, modified_files, unchanged_files


# ============================================================
# MAIN
# ============================================================

def main():
    full_mode = "--full" in sys.argv
    list_mode = "--list" in sys.argv
    status_mode = "--status" in sys.argv

    # Single-instance PRIMA di scan e load_registry: due run concorrenti si
    # sovrascrivono il registry a vicenda (read-modify-write non atomico) e
    # possono interlacciare delete e upsert sugli stessi chunk_id.
    # --list/--status non scrivono nulla: restano disponibili durante un run.
    # NB: lock_fd va tenuto referenziato per tutta la durata di main() — se il
    # file object viene raccolto, l'OS rilascia il lock. Non "pulire" come
    # variabile inutilizzata.
    lock_fd = None
    if not (list_mode or status_mode):
        lock_fd = acquire_single_instance_lock()
        if lock_fd is None:
            print(f"Un'altra istanza e' gia' in esecuzione ({LOCK_FILE}). Esco.")
            sys.exit(1)

    print(f"TAILOR Document Ingest")
    print(f"Cartelle monitorate: {len(WATCH_FOLDERS)}")
    for folder in WATCH_FOLDERS:
        exists = "OK" if os.path.exists(folder) else "NON TROVATA"
        print(f"  [{exists}] {folder}")

    # Scan
    print(f"\nScansione file...")
    files = scan_folders()
    print(f"  Trovati {len(files)} file supportati")

    if not files:
        print("No files to process.")
        return

    # List mode: mostra e esci
    if list_mode:
        print(f"\nFile trovati:")
        for f in sorted(files, key=lambda x: x["filepath"]):
            print(f"  [{f['extension']}] {f['filename']} ({f['size_kb']:.0f} KB) - {f['modified_date']}")
            print(f"         {f['filepath']}")
        return

    # Classifica
    registry = load_registry()
    new_files, modified_files, unchanged_files = classify_files(files, registry, full_mode)

    if status_mode:
        print(f"\nStato file:")
        print(f"  Nuovi:      {len(new_files)}")
        print(f"  Modificati: {len(modified_files)}")
        print(f"  Unchanged:  {len(unchanged_files)}")
        for f in new_files:
            print(f"    [NUOVO]      {f['filename']}")
        for f in modified_files:
            print(f"    [MODIFICATO] {f['filename']}")
        return

    to_process = new_files + modified_files
    if not to_process:
        print(f"\nAll files already in KB ({len(unchanged_files)} unchanged). Nothing to do.")
        return

    mode_label = "FULL RE-INGEST" if full_mode else "INCREMENTALE"
    print(f"\nMode: {mode_label}")
    print(f"  Nuovi:      {len(new_files)}")
    print(f"  Modificati: {len(modified_files)}")
    print(f"  Unchanged:  {len(unchanged_files)} (skipped)")

    # Setup
    os.makedirs(DB_DIR, exist_ok=True)
    client = chromadb.PersistentClient(path=DB_DIR)
    collection = client.get_collection(name=COLLECTION_NAME)

    total_chunks_added = 0
    total_errors = 0
    start_time = time.time()
    run_id = make_run_id("document")

    for idx, f in enumerate(to_process):
        print(f"\n[{idx+1}/{len(to_process)}] {f['filename']} ({f['size_kb']:.0f} KB)")

        # Estrazione PRIMA di toccare la KB: un documento sparisce solo se
        # sostituito, mai perche' l'estrazione e' fallita.
        pre_doctype = infer_doc_type(f["filename"])
        sections = extract_text(f["filepath"], doctype=pre_doctype)
        if not sections:
            # Nessun cleanup: i chunk vecchi e i loro sidecar restano intatti.
            fail_count = mark_file_failed(registry, f, "no_text_extracted")
            print(f"  No text extracted, skip (fail_count={fail_count})")
            continue

        total_chars = sum(len(s["text"]) for s in sections)
        print(f"  Estratte {len(sections)} sezioni ({total_chars:,} chars)")

        # Chunk
        chunks = chunk_document(sections, f)
        print(f"  Generati {len(chunks)} chunk")

        # Embedding PRIMA del cleanup: e' la fase lunga e quella che fallisce
        # piu' spesso (rete, provider). Facendola qui la finestra in cui il
        # documento e' assente dalla KB si riduce alla sola coppia
        # delete+upsert. In RAM sta un file alla volta: dimensioni banali.
        batches, embed_errors = embed_chunks(chunks)
        if embed_errors or not batches:
            # KB non ancora toccata: il documento vecchio resta al suo posto.
            total_errors += embed_errors or 1
            fail_count = mark_file_failed(registry, f, "embedding_failed")
            print(f"  Embedding fallito, KB intatta, file SALTATO "
                  f"(fail_count={fail_count})")
            continue

        # Da qui la sequenza e' atomica per-file: la delete deve comunque
        # precedere l'upsert, altrimenti si scriverebbe sopra chunk vecchi
        # ancora presenti. Se il cleanup fallisce il file viene saltato.
        try:
            _pre_upsert_cleanup(collection, f, registry)
        except Exception as e:
            print(f"  ERRORE cleanup pre-upsert, file SALTATO: {e}")
            total_errors += 1
            mark_file_failed(registry, f, f"cleanup_failed: {e}")
            continue

        # Ingest
        processed, errors = ingest_chunks(collection, batches, run_id=run_id)
        total_chunks_added += processed
        total_errors += errors
        print(f"  Ingestati {processed} chunk" + (f" ({errors} errori)" if errors else ""))

        # Registry: successo SOLO se il documento e' integro. Un upsert
        # parziale registrato come riuscito renderebbe il file "unchanged" al
        # giro dopo, cristallizzando in KB un documento monco.
        if errors == 0 and processed == len(chunks):
            registry[f["filepath"]] = {
                "hash": f["hash"],
                "filename": f["filename"],
                "date_ingested": time.strftime("%Y-%m-%d %H:%M:%S"),
                "chunks": len(chunks),
                "chars": total_chars,
                "sections": len(sections),
            }
        else:
            fail_count = mark_file_failed(
                registry, f, f"partial_upsert: {processed}/{len(chunks)} chunk")
            print(f"  ATTENZIONE upsert parziale ({processed}/{len(chunks)}): "
                  f"registrato come FALLITO, retry al giro dopo "
                  f"(fail_count={fail_count})")

    # Salva registry
    save_registry(registry)

    elapsed = time.time() - start_time
    print(f"\n{'='*50}")
    print(f"Result:")
    print(f"  File processati: {len(to_process)}")
    print(f"  Chunk aggiunti: {total_chunks_added:,}")
    print(f"  Errori: {total_errors}")
    print(f"  Tempo: {elapsed:.1f}s")
    print(f"  Totale chunk in KB: {collection.count():,}")

    # Fail-loud: batch falliti -> exit non-zero (il cron non deve tacere).
    if total_errors > 0:
        sys.exit(1)


if __name__ == "__main__":
    main()